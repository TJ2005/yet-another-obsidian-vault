"""Maximum-data document extraction -> Obsidian markdown.

Single entry point for the parse-docs skill. Handles PDF, PPTX, DOCX, XLSX,
RTF/.doc/.odt/HTML (via macOS textutil), CSV, and images (tesseract OCR).
"""

# /// script
# requires-python = ">=3.11"
# dependencies = [
#   "pymupdf>=1.24",
#   "python-pptx>=0.6.23",
#   "python-docx>=1.1",
#   "openpyxl>=3.1",
#   "mammoth>=1.6",
#   "html2text>=2020.1.16",
#   "beautifulsoup4>=4.12",
# ]
# ///

from __future__ import annotations

import argparse
import base64
import csv
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import date
from pathlib import Path
from typing import Optional

SUPPORTED = {
    ".pdf", ".pptx", ".docx", ".xlsx", ".rtf", ".doc", ".odt", ".html", ".htm",
    ".csv", ".txt", ".md", ".markdown", ".png", ".jpg", ".jpeg", ".tiff", ".tif",
    ".webp", ".bmp", ".gif",
}
LEGACY = {".ppt", ".xls", ".odp", ".ods", ".msg", ".pages", ".key"}


def out_name(path: Path) -> str:
    return f"{path.stem}.extracted.md"


def assets_dir(outdir: Path, path: Path) -> Path:
    d = outdir / f"{path.stem}.extracted_assets"
    d.mkdir(parents=True, exist_ok=True)
    return d


def fm(path: Path, extra: Optional[str] = None) -> str:
    head = (
        "---\n"
        f"source: {path.name}\n"
        f"source_type: {path.suffix.lstrip('.').lower()}\n"
        f"extracted: {date.today().isoformat()}\n"
    )
    if extra:
        head += extra
    head += "---\n\n"
    return head


class AssetWriter:
    def __init__(self, adir: Path):
        self.adir = adir
        self.count = 0
        self._seen: set = set()

    def add(self, data: bytes, ext: str, label: str = "img", root="") -> str:
        self.count += 1
        name = f"{root}{label}_{self.count}.{ext}"
        (self.adir / name).write_bytes(data)
        return name

    def add_unique(self, data: bytes, ext: str, label: str = "img", root="") -> Optional[str]:
        key = (label, len(data))
        if key in self._seen:
            return None
        self._seen.add(key)
        return self.add(data, ext, label, root)


def wikilink(adir: Path, name: str) -> str:
    return f"![[{adir.name}/{name}]]"


# ---------------------------------------------------------------- textutil
def textutil(path: Path) -> Optional[str]:
    r = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        capture_output=True, text=True,
    )
    if r.returncode != 0:
        return None
    return r.stdout.strip()


# ------------------------------------------------------------------- PDF
def extract_pdf(path: Path, aw: AssetWriter, ocr_all: bool, lang: str) -> str:
    import fitz

    doc = fitz.open(path)
    if doc.needs_pass:
        lock = doc.metadata.get("encryption") or "password-protected"
        return f"> Cannot extract: PDF is {lock}. Ask the user for the password.\n"

    parts = []
    meta = {k: v for k, v in doc.metadata.items() if v}
    if meta:
        parts.append("## Metadata\n")
        for k, v in meta.items():
            parts.append(f"- **{k}**: {v}")
        parts.append("\n")

    toc = doc.get_toc()
    if toc:
        parts.append("## Table of Contents\n")
        for lvl, title, page in toc:
            parts.append(f"{'  ' * (lvl - 1)}- {title} (p. {page})")
        parts.append("\n")

    ocr_count = 0
    for i, page in enumerate(doc, 1):
        parts.append(f"\n## Page {i}\n")
        md = _page_md(page).strip()

        images = page.get_images(full=True)
        imgs_here = []
        for img in images:
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                ext = "png" if pix.n - pix.alpha < 4 else "jpg"
                name = aw.add_unique(pix.tobytes(ext), ext, root=f"p{i}_")
                if name:
                    imgs_here.append(wikilink(aw.adir, name))
            except Exception:
                pass

        if (not md or ocr_all) and (imgs_here or page.get_images(full=True)):
            ocr_text = _ocr_page(page, lang)
            if ocr_text.strip():
                ocr_count += 1
                md = f"[OCR]\n\n{ocr_text.strip()}"
        elif not md and not ocr_all:
            parts.append("_No text or images found on this page._\n")

        if imgs_here:
            parts.append("**Images:** " + " ".join(imgs_here) + "\n")
        if md:
            parts.append(md + "\n")

        links = [l for l in page.get_links() if l.get("uri")]
        if links:
            parts.append("**Links:**\n")
            for l in links:
                parts.append(f"- [{l['uri']}]({l['uri']})")
            parts.append("\n")

    extra = f"pages: {doc.page_count}\nocr_pages: {ocr_count}\n"
    doc.close()
    return fm(path, extra) + "\n".join(parts)


def _page_md(page) -> str:
    for fn_name in ("get_markdown", "get_text_markdown"):
        fn = getattr(page, fn_name, None)
        if fn:
            try:
                return fn()
            except Exception:
                pass
    try:
        return page.get_text("markdown")
    except Exception:
        return page.get_text("text")


def _ocr_page(page, lang: str) -> str:
    import fitz

    if not shutil.which("tesseract"):
        return ""
    pix = page.get_pixmap(matrix=fitz.Matrix(3, 3), colorspace=fitz.csRGB)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
        tmp = tf.name
    try:
        pix.save(tmp)
        r = subprocess.run(
            ["tesseract", tmp, "stdout", "-l", lang],
            capture_output=True, text=True,
        )
        return r.stdout if r.returncode == 0 else ""
    finally:
        Path(tmp).unlink(missing_ok=True)


# ------------------------------------------------------------------ PPTX
def extract_pptx(path: Path, aw: AssetWriter) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    parts = [fm(path, f"slides: {len(prs.slides)}\n")]
    n_img = 0

    def walk(shapes, level: int, slide_no: int, out: list[str]) -> None:
        nonlocal n_img
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    walk(shape.shapes, level + 1, slide_no, out)
                except Exception:
                    pass
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        continue
                    indent = "  " * (level + para.level)
                    bullets = "  - " if (para.level > 0 or para.text.strip().startswith(("•", "-", "1.", "*"))) else ""
                    if para.level == 0 and level == 0 and shape == slide_title_ref:
                        out.append(f"# {text}")
                    else:
                        out.append(f"{indent}{bullets}{text}")
            if shape.has_table:
                tbl = shape.table
                out.append("\n")
                for rrow in tbl.rows:
                    cells = [c.text.replace("\n", " ").strip() for c in rrow.cells]
                    out.append("| " + " | ".join(cells) + " |")
                    if rrow is tbl.rows[0]:
                        out.append("| " + " | ".join(["---"] * len(cells)) + " |")
                out.append("\n")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    n_img += 1
                    name = aw.add(image.blob, image.ext or "png", root=f"s{slide_no}_")
                    out.append(wikilink(aw.adir, name))
                except Exception:
                    pass
            if getattr(shape, "has_chart", False):
                try:
                    chart = shape.chart
                    sname = " / ".join(series.name for series in chart.plots[0].series) if chart.plots else "chart"
                    out.append(f"\n**Chart: {sname}**\n")
                    out.append("| Label | Value |")
                    out.append("| --- | --- |")
                    for series in chart.plots[0].series:
                        cats = series.categories if hasattr(series, "categories") else None
                        vals = list(series.values)
                        for i, v in enumerate(vals):
                            c = str(cats[i]) if cats and i < len(cats) else str(i + 1)
                            out.append(f"| {c} | {v} |")
                    out.append("\n")
                except Exception:
                    pass

    for si, slide in enumerate(prs.slides, 1):
        slide_title_ref = None
        slide_out: list[str] = []
        try:
            if slide.shapes.title is not None:
                slide_title_ref = slide.shapes.title
        except Exception:
            pass
        slide_out.append(f"## Slide {si}\n")
        try:
            walk(slide.shapes, 0, si, slide_out)
        except Exception as e:
            slide_out.append(f"_Error walking slide: {e}_")
        notes = None
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        if notes:
            slide_out.append(f"\n### Speaker notes\n\n{notes}\n")
        parts.append("\n".join(slide_out))

    extra = f"images: {n_img}\n"
    return fm(path, extra) + "\n".join(parts)


# ------------------------------------------------------------------ DOCX
def extract_docx(path: Path, aw: AssetWriter) -> str:
    import mammoth

    with open(path, "rb") as f:
        result = mammoth.convert_to_html(f)
    md = html2markdown(result.value or f"*Extraction failed: {result.messages}*")

    md = _inline_images(md, aw)
    return fm(path) + md


def _inline_images(md: str, aw: AssetWriter) -> str:
    pat = re.compile(r"!\[[^\]]*\]\(data:image/(?P<ext>png|jpe?g|gif|webp|bmp);base64,(?P<b64>[^)]+)\)", re.I)

    def repl(m: re.Match) -> str:
        try:
            data = base64.b64decode(m.group("b64"))
            ext = "jpg" if m.group("ext").startswith("jpeg") or m.group("ext") == "jpg" else m.group("ext").lower()
            name = aw.add(data, ext)
            return wikilink(aw.adir, name)
        except Exception:
            return m.group(0)

    return pat.sub(repl, md)


def html2markdown(html_text: str) -> str:
    import html2text

    h = html2text.HTML2Text()
    h.body_width = 0
    h.ignore_links = False
    h.ignore_images = False
    h.single_line_break = False
    return h.handle(html_text).strip()


# ------------------------------------------------------------------ XLSX
def extract_xlsx(path: Path, aw: AssetWriter) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True)
    wbf = load_workbook(path, data_only=False)
    parts = [fm(path, f"sheets: {len(wb.worksheets)}\n")]

    for ws, wsf in zip(wb.worksheets, wbf.worksheets):
        parts.append(f"## Sheet: {ws.title}\n")
        merged = sorted(str(r) for r in ws.merged_cells.ranges)

        rows_out: list[str] = []
        for rrow in ws.iter_rows():
            vals: list[str] = []
            for cell in rrow:
                v = cell.value
                fv = wsf.cell(row=cell.row, column=cell.column)
                if fv.data_type == "f":
                    v = f"{fv.value}  `{v}`" if v is not None else f"`{fv.value}`"
                if v is None:
                    vals.append("")
                elif isinstance(v, (list, tuple, dict)):
                    vals.append(str(v))
                else:
                    vals.append(str(v).replace("|", "\\|").replace("\n", " "))
            if any(vals):
                rows_out.append("| " + " | ".join(vals) + " |")

        if rows_out:
            width = max(len(r.split("|")) - 2 for r in rows_out)
            parts.append("\n".join(rows_out[:1]))
            parts.append("| " + " | ".join(["---"] * width) + " |")
            parts.extend(rows_out[1:])
        else:
            parts.append("_Empty sheet._\n")

        if merged:
            parts.append("\n**Merged cells:** " + ", ".join(merged) + "\n")

        comments = []
        for row in ws.iter_rows():
            for cell in row:
                if cell.comment and cell.comment.text.strip():
                    comments.append(f"- **{cell.coordinate}**: {cell.comment.text.strip()}")
        if comments:
            parts.append("\n**Comments:**\n" + "\n".join(comments) + "\n")

        try:
            for img in ws._images:
                data = img._data() if hasattr(img, "_data") else None
                if data:
                    fmt = getattr(img, "format", None) or "png"
                    ext = str(fmt).lower() if isinstance(fmt, str) else fmt
                    if ext in ("jpeg", "jpg"):
                        ext = "jpg"
                    name = aw.add(data, ext, root=f"{ws.title}_")
                    parts.append(wikilink(aw.adir, name))
        except Exception:
            pass

    return "\n".join(parts)


# ------------------------------------------------------------------- CSV
def extract_csv(path: Path) -> str:
    with open(path, newline="", encoding="utf-8-sig", errors="replace") as f:
        rows = list(csv.reader(f))
    if not rows:
        return fm(path) + "_Empty file._\n"
    width = max(len(r) for r in rows)
    for r in rows:
        while len(r) < width:
            r.append("")
    out = [fm(path), "| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    out.extend("| " + " | ".join(r) + " |" for r in rows[1:])
    return "\n".join(out)


# ----------------------------------------------------------------- image
def extract_image(path: Path, lang: str, aw: AssetWriter) -> str:
    if not shutil.which("tesseract"):
        return fm(path) + "_OCR unavailable: `brew install tesseract`._\n"
    r = subprocess.run([shutil.which("tesseract"), str(path), "stdout", "-l", lang], capture_output=True, text=True)
    text = r.stdout.strip() if r.returncode == 0 else ""
    dest = aw.adir / path.name
    shutil.copy2(path, dest)
    body = f"## Source image\n\n{wikilink(aw.adir, dest.name)}\n\n"
    if text:
        body += f"[OCR]\n\n{text}"
    else:
        body += "_No text recognized._"
    return fm(path) + body


# ------------------------------------------------------------------ main
def process(path: Path, outdir: Path, ocr_all: bool, lang: str) -> Optional[str]:
    ext = path.suffix.lower()
    aw = AssetWriter(assets_dir(outdir, path))

    if ext == ".pdf":
        content = extract_pdf(path, aw, ocr_all, lang)
    elif ext == ".pptx":
        content = extract_pptx(path, aw)
    elif ext == ".docx":
        content = extract_docx(path, aw)
    elif ext == ".xlsx":
        content = extract_xlsx(path, aw)
    elif ext == ".csv":
        content = extract_csv(path)
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".webp", ".bmp", ".gif"):
        content = extract_image(path, lang, aw)
    elif ext in (".rtf", ".doc", ".odt", ".html", ".htm", ".txt", ".md", ".markdown"):
        txt = textutil(path) if ext not in (".txt", ".md", ".markdown") else path.read_text(errors="replace")
        content = fm(path) + (txt.strip() if txt else "_Empty file._\n")
    elif ext in LEGACY:
        return f"LEGACY {path.name}: unsupported format {ext} — convert with LibreOffice or Google Drive first."
    else:
        return f"SKIP {path.name}: unhandled extension {ext}"

    outp = outdir / out_name(path)
    outp.write_text(content, encoding="utf-8")
    n_img = aw.count
    return f"{'OK' if outp.exists() else 'FAIL'} {path.name} -> {outp.name}" + (f" ({n_img} asset(s))" if n_img else "")


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract maximum data from documents -> Obsidian markdown")
    ap.add_argument("paths", nargs="+", help="Files or directory to ingest")
    ap.add_argument("--out", default="Extracted", help="Output directory (default: ./Extracted)")
    ap.add_argument("--ocr-all", action="store_true", help="OCR every PDF page even with a text layer")
    ap.add_argument("--recursive", action="store_true", help="Recurse into subdirectories")
    ap.add_argument("--lang", default="eng", help="Tesseract language(s), e.g. eng, eng+spa")
    args = ap.parse_args()

    outdir = Path(args.out)
    outdir.mkdir(parents=True, exist_ok=True)

    files: list[Path] = []
    for p in args.paths:
        pp = Path(p)
        if pp.is_dir():
            files.extend(sorted(f for f in (pp.rglob("*") if args.recursive else pp.glob("*")) if f.suffix.lower() in SUPPORTED))
        elif pp.suffix.lower() in SUPPORTED:
            files.append(pp)
        else:
            print(f"SKIP {pp.name}: unsupported extension {pp.suffix}", file=sys.stderr)

    if not files:
        print("No supported files found.", file=sys.stderr)
        return 1

    ok = 0
    for f in files:
        try:
            res = process(f, outdir, args.ocr_all, args.lang)
            print(res or "", file=sys.stderr)
            if res and res.startswith("OK"):
                ok += 1
        except Exception as e:
            print(f"ERR {f.name}: {e}", file=sys.stderr)

    print(f"Done: {ok}/{len(files)} extracted -> {outdir}/", file=sys.stderr)
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())